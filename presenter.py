# Is called by:
# - datachecker
# - tcm

import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches


class Presenter(object):
    def __init__(self, new_pres=None, existing_pres=None, title='', subtitle=''):
        self.prs = Presentation(existing_pres)
        self.name = existing_pres

        if new_pres is not None:
            self.name = new_pres

            self.start_presentation(title=title, subtitle=subtitle)
            self.save_presentation(new_pres)

    def start_presentation(self, title='', subtitle=''):
        self.title = title
        self.subtitle = subtitle

        # Title (presentation title slide)
        slide_layout = self.prs.slide_layouts[0]
        top_slide = self.prs.slides.add_slide(slide_layout)

        top_slide.placeholders[0].text = title
        top_slide.placeholders[1].text = subtitle

    def add_picture_slide(self, picture, title=''):
        # Title only slide
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        slide.placeholders[0].text = title
        left = Inches(1)
        top = Inches(1.3)
        slide.shapes.add_picture(picture, left, top)

    def add_text_slide(self, text, title=''):
        # Title and content slide
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        slide.placeholders[0].text = title
        slide.placeholders[1].text = text

    def add_graph_slide(self, x=None, y=None, graph_type='Plot', **kwargs):
        # Default values
        labels = None
        legend_loc = 'upper left'
        slide_title = ''

        fig = plt.figure()

        if 'xticks' in kwargs:
            plt.xticks(rotation=kwargs['xticks'])
        if 'labels' in kwargs:
            labels = kwargs['labels']
        if 'legend_loc' in kwargs:
            legend_loc = kwargs['legend_loc']

        # The actual plotting
        if graph_type is 'scatter':
            plt.scatter(x, y)
        elif graph_type is 'plot':
            if np.ndim(y) > 1:
                if labels is None:
                    labels = np.array2string(np.arange(len(y)))
                for i in range(0, np.shape(y)[0]):
                    plt.plot(x, y[i], label=labels[i])
            else:
                if labels is None:
                    plt.plot(x, y)
                else:
                    plt.plot(x, y, label=labels)
            if labels is not None:
                plt.legend(loc=legend_loc)

        # Set options for graph
        if 'slide_title' in kwargs:
            slide_title = kwargs['slide_title']
        if 'xlabel' in kwargs:
            plt.xlabel(kwargs['xlabel'])
        if 'ylabel' in kwargs:
            plt.ylabel(kwargs['ylabel'])
        if 'plot_title' in kwargs:
            plt.title(kwargs['plot_title'])
        if 'grid' in kwargs:
            plt.grid(kwargs['grid'])
        if 'xlim' in kwargs:
            if kwargs['xlim'] is None:
                plt.autoscale(True, axis='x')
            else:
                plt.xlim(kwargs['xlim'])
        if 'ylim' in kwargs:
            if kwargs['ylim'] is None:
                plt.autoscale(True, axis='y')
            else:
                plt.ylim(kwargs['ylim'])

        img_path = 'temp_PRESENTER.png'
        plt.savefig(img_path)
        plt.close(fig)

        # Adding a slide
        self.add_picture_slide(img_path, title=slide_title)

    def save_presentation(self, path=None):
        if path is None:
            self.prs.save(self.name)
        else:
            self.prs.save(path)
